import streamlit as st
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz  # PyMuPDF
import io
import os
import zipfile
import re # For regex and splitting

# --- GLOBAL CONSTANTS ---
COMMON_HONORIFICS = [
    "Mr.", "Mrs.", "Ms.", "Miss", "Dr.", "Doctor", "Prof.", "Professor",
    "Rev.", "Reverend", "Hon.", "Honorable", "Gen.", "General", "Sgt.", "Sergeant",
    "Capt.", "Captain", "Lt.", "Lieutenant", "Fr.", "Father", "Sr.", "Sister"
]

# --- HELPER FUNCTIONS ---

def parse_names_from_ui(names_input_str_ui):
    """Parses the raw comma-separated string from the UI."""
    if not names_input_str_ui: return []
    names = [name.strip() for name in names_input_str_ui.split(',') if name.strip()]
    return names

def generate_name_variations(user_provided_names_list, honorifics_list):
    """
    Generates a comprehensive list of name variations for redaction.
    Input: A list of names as directly provided by the user.
    """
    if not user_provided_names_list:
        return []

    variations = set()
    normalized_honorifics_set = set(h.lower().replace('.', '') for h in honorifics_list)

    for name_entry in user_provided_names_list:
        original_name_entry = name_entry.strip()
        if not original_name_entry: continue
        variations.add(original_name_entry)

        # Split the entry into parts for analysis (handles spaces, hyphens, periods)
        name_parts_for_logic = re.split(r'[\s.-]+', original_name_entry)
        name_parts_for_logic = [p for p in name_parts_for_logic if p] # Clean empty parts

        name_without_honorific = original_name_entry # Default
        if name_parts_for_logic:
            potential_honorific_in_input = name_parts_for_logic[0].lower().replace('.', '')
            if potential_honorific_in_input in normalized_honorifics_set:
                if len(name_parts_for_logic) > 1:
                    # If "Prof. John Doe", name_without_honorific becomes "John Doe"
                    name_without_honorific = " ".join(name_parts_for_logic[1:])
                    variations.add(name_without_honorific) # Add "John Doe"
            # If no honorific found at the start, name_without_honorific remains original_name_entry

        # Now, use 'name_without_honorific' to generate more variations
        # 1. Add variations WITH all common honorifics
        for honorific in honorifics_list: # Use original honorifics list to keep their format
            variations.add(f"{honorific} {name_without_honorific}")
            # Also add honorific + just the last part if name_without_honorific has multiple words
            base_name_parts_for_honorific = name_without_honorific.split() # Simple space split
            if len(base_name_parts_for_honorific) > 1:
                 variations.add(f"{honorific} {base_name_parts_for_honorific[-1]}") # e.g., "Dr. Doe" from "John Doe"

        # 2. Handle initials and multi-part names from 'name_without_honorific'
        base_parts = name_without_honorific.split() # Simple space split for this logic
        
        if len(base_parts) == 1: # Single word like "Young"
            variations.add(base_parts[0]) # Already added via name_without_honorific usually
        
        elif len(base_parts) == 2: # e.g., "John Doe"
            first, last = base_parts[0], base_parts[1]
            variations.add(last) # "Doe"
            if first and last: # Ensure parts are not empty
                if len(first)>0 : variations.add(f"{first[0]}. {last}") # "J. Doe"
                if len(last)>0 : variations.add(f"{first} {last[0]}.") # "John D."
                if len(first)>0 and len(last)>0 : variations.add(f"{first[0]}.{last[0]}.") # "J.D."

        elif len(base_parts) == 3: # e.g., "John King Doe" or "John K Doe" (if K has no period)
            first, middle, last = base_parts[0], base_parts[1], base_parts[2]
            variations.add(last) # "Doe"
            variations.add(f"{first} {last}") # "John Doe"
            if first and middle and last:
                if len(first)>0: variations.add(f"{first[0]}. {last}") # "J. Doe" (ignoring middle)
                if len(first)>0 and len(middle)>0: variations.add(f"{first[0]}. {middle[0]}. {last}") # "J. K. Doe"
                if len(first)>0 and len(middle)>0: variations.add(f"{first} {middle[0]}. {last}") # "John K. Doe"

    # Final sort and filter: remove very short items (e.g. length 1) unless specifically intended.
    # For now, min length 2 seems reasonable to avoid redacting "A" if "A. B. Cee" was processed.
    # This also filters out empty strings that might have crept in.
    sorted_variations = sorted(list(v for v in variations if v and len(v.strip()) > 1), key=len, reverse=True)
    return sorted_variations


def redact_text_in_runs(runs, names_to_redact_variations, redaction_string="[REDACTED]"):
    """Iterates through runs, uses regex for whole-word, case-insensitive replacement."""
    modified_in_paragraph = False
    for run in runs:
        current_run_text = run.text
        if not current_run_text.strip(): # Skip empty or whitespace-only runs
            continue
        
        text_modified_in_this_run = False
        for name_var in names_to_redact_variations:
            if not name_var: continue # Skip empty variations, should be filtered by generate_name_variations

            # re.escape to handle special characters in names (e.g., "Dr. No")
            # \b for word boundaries
            # Using try-except for re.subn as a safeguard
            try:
                # Pattern needs to be re-evaluated if current_run_text changed
                pattern = r'\b' + re.escape(name_var) + r'\b'
                new_text, num_replacements = re.subn(pattern, redaction_string, current_run_text, flags=re.IGNORECASE)
                if num_replacements > 0:
                    current_run_text = new_text # Update text for next variation in this run
                    text_modified_in_this_run = True
            except re.error as e:
                # This should be rare with re.escape, but good to have.
                st.warning(f"Regex error for name variation '{name_var}': {e}. Skipping this variation for this run segment.")
                continue
        
        if text_modified_in_this_run:
            run.text = current_run_text # Apply accumulated changes to the run object
            modified_in_paragraph = True
            
    return modified_in_paragraph

def redact_docx(docx_file_stream, names_variations_list, redaction_string):
    doc = Document(docx_file_stream)
    modified_doc = False
    for para in doc.paragraphs:
        if redact_text_in_runs(para.runs, names_variations_list, redaction_string): modified_doc = True
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if redact_text_in_runs(para.runs, names_variations_list, redaction_string): modified_doc = True
    for section in doc.sections:
        for header_para in section.header.paragraphs:
            if redact_text_in_runs(header_para.runs, names_variations_list, redaction_string): modified_doc = True
        for footer_para in section.footer.paragraphs:
            if redact_text_in_runs(footer_para.runs, names_variations_list, redaction_string): modified_doc = True
    if modified_doc:
        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        return bio
    return None

def redact_pptx(pptx_file_stream, names_variations_list, redaction_string):
    prs = Presentation(pptx_file_stream)
    modified_prs = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if redact_text_in_runs(para.runs, names_variations_list, redaction_string): modified_prs = True
            if shape.has_table:
                table = shape.table
                for r_idx in range(len(table.rows)):
                    for c_idx in range(len(table.columns)):
                        cell = table.cell(r_idx, c_idx)
                        if cell.text_frame:
                            for para in cell.text_frame.paragraphs:
                                if redact_text_in_runs(para.runs, names_variations_list, redaction_string): modified_prs = True
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                 if redact_text_in_runs(para.runs, names_variations_list, redaction_string): modified_prs = True
    if modified_prs:
        bio = io.BytesIO()
        prs.save(bio)
        bio.seek(0)
        return bio
    return None

def redact_pdf(pdf_file_stream, names_to_redact_variations, redaction_string): # redaction_string not used for PDF visual
    pdf_bytes = pdf_file_stream.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    modified_pdf = False
    
    # PyMuPDF flags: TEXT_PRESERVE_LIGATURES=1, TEXT_PRESERVE_WHITESPACE=2
    # search_for is case-insensitive by default in recent versions (>=1.18.0).
    search_flags = fitz.TEXT_PRESERVE_LIGATURES | fitz.TEXT_PRESERVE_WHITESPACE 

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        all_rects_to_redact_on_page = [] # Collect all rects for this page

        for name_var in names_to_redact_variations:
            if not name_var: continue # Should be filtered by generate_name_variations

            try:
                # search_for returns a list of Rect objects (or Quads if quads=True)
                text_instances = page.search_for(name_var, flags=search_flags, quads=False) 
                for inst_rect in text_instances:
                    if not inst_rect.is_empty and not inst_rect.is_infinite: # Validate rect
                        all_rects_to_redact_on_page.append(inst_rect)
            except AttributeError: # Fallback for older PyMuPDF that might not have flags structured the same
                text_instances = page.search_for(name_var, quads=False) 
                for inst_rect in text_instances:
                    if not inst_rect.is_empty and not inst_rect.is_infinite:
                        all_rects_to_redact_on_page.append(inst_rect)
            except Exception as e: # Catch other potential search errors
                st.warning(f"PDF search error for '{name_var}' on page {page_num + 1}: {e}")

        if all_rects_to_redact_on_page:
            modified_pdf = True 
            for r in all_rects_to_redact_on_page:
                try:
                    page.add_redact_annot(r, text="", fill=(0, 0, 0)) # Add redaction annotation
                except Exception as annot_e:
                    # st.warning(f"Could not add PDF redaction annotation for rect {r} on page {page_num+1}: {annot_e}")
                    pass # Continue if one annotation fails, to not halt all page redaction
            try:
                page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE) # Apply all on page
            except Exception as apply_e:
                # st.error(f"Error applying PDF redactions on page {page_num+1}: {apply_e}")
                pass
    
    if modified_pdf:
        bio = io.BytesIO()
        doc.save(bio, garbage=3, deflate=True, clean=True) # Save efficiently
        bio.seek(0)
        doc.close()
        return bio
    doc.close()
    return None

def extract_initials_from_filename(filename_base):
    """Attempts to extract 2 initials from a filename base."""
    # Remove file extension if present (though filename_base should ideally be pre-stripped)
    name_part = os.path.splitext(filename_base)[0]
    
    # Split by common delimiters: space, hyphen, underscore
    parts = re.split(r'[\s_-]+', name_part)
    parts = [p for p in parts if p and p[0].isalpha()] # Keep only parts starting with a letter

    initials = ""
    if len(parts) >= 2: # Typically "FirstName LastName ..." or "Word1 Word2 ..."
        initials = parts[0][0] + parts[1][0]
    elif len(parts) == 1 and len(parts[0]) >=2 : # Single word, take first two letters if long enough
        initials = parts[0][:2]
    elif len(parts) == 1 and len(parts[0]) == 1: # Single letter word
        initials = parts[0][0]
        
    return initials.upper() if initials else ""


# --- Streamlit App UI and Main Logic ---
st.set_page_config(layout="wide", page_title="TU Name Redactor", page_icon="🐉")

st.markdown("""
<style>
    .stButton>button[kind="primary"] { 
        background-color: #4CAF50 !important; color: white !important; border: none !important;
    }
    .stButton>button[kind="primary"]:hover {
        background-color: #45a049 !important; color: white !important;
    }
    .stButton>button[kind="primary"]:active {
        background-color: #3e8e41 !important; color: white !important;
    }
</style>
""", unsafe_allow_html=True)

st.title("🐉 TU Name Redactor")
st.markdown("Welcome! This tool redacts names from documents. Configure settings, upload files, and download.")
st.markdown("---")

with st.sidebar:
    st.header("⚙️ Redaction Settings")
    names_input_str_ui = st.text_area( 
        "Names to redact (comma-separated):", 
        placeholder="e.g., Professor Plum, John K. Doe, Ms. Scarlet, Young",
        height=120, 
        help="Enter full names, names with honorifics, or last names. The app will generate variations."
    )
    redaction_text_docx_pptx = st.text_input(
        "Redaction text for DOCX/PPTX:", value="[REDACTED]",
        help="Text replacement for Word/PowerPoint. PDFs are blacked out.")
    st.markdown("---")
    st.subheader("📦 ZIP Output Options")
    st.caption("(These options apply *only* when processing an uploaded ZIP file)")
    output_zip_name_user_input = st.text_input(
        "Custom output ZIP name base:", 
        placeholder="e.g., MyProject_Redacted",
        help="Base name for the output ZIP. Defaults to 'redacted_[original_zip_name]'. '.zip' added automatically."
    )
    st.caption("Initials for files within the ZIP are auto-extracted from their original filenames (if possible, max 2 initials).")

uploaded_files = st.file_uploader(
    "Upload .docx, .pptx, .pdf, or .zip files", type=["docx", "pptx", "pdf", "zip"],
    accept_multiple_files=True, help="Upload individual files or ZIP archives.")

if st.button("Redact Files", type="primary", use_container_width=True, key="redact_button"):
    if not uploaded_files: st.warning("⚠️ Please upload at least one file.")
    elif not names_input_str_ui.strip(): st.warning("⚠️ Please enter at least one name to redact.")
    else:
        # 1. Get raw names from UI
        parsed_user_names = parse_names_from_ui(names_input_str_ui)
        # 2. Generate variations
        names_variations_list = generate_name_variations(parsed_user_names, COMMON_HONORIFICS)
        
        if not names_variations_list:
            st.warning("⚠️ No valid names or variations generated for redaction. Please check your input.")
        else:
            st.info(f"Attempting to redact based on {len(names_variations_list)} name variations (e.g., {', '.join(names_variations_list[:min(3, len(names_variations_list))])}...).")
            st.info(f"DOCX/PPTX redaction: '{redaction_text_docx_pptx}'. PDFs blacked out.")
            
            overall_docs_modified_count = 0
            files_for_download = [] 

            with st.spinner("🔧 Processing files... This might take a moment..."):
                for uploaded_file_obj in uploaded_files:
                    original_input_name = uploaded_file_obj.name
                    file_extension = os.path.splitext(original_input_name)[1].lower()
                    
                    if file_extension == ".zip":
                        st.write(f"--- Processing ZIP: **{original_input_name}** ---")
                        processed_zip_members_data = [] 
                        members_found_for_processing_in_zip = False
                        custom_zip_name_base_from_input = output_zip_name_user_input.strip()
                        actual_output_zip_base_name = custom_zip_name_base_from_input if custom_zip_name_base_from_input \
                                                      else f"redacted_{os.path.splitext(original_input_name)[0]}"
                        try:
                            with zipfile.ZipFile(uploaded_file_obj, 'r') as zip_ref:
                                for member_name_in_zip in zip_ref.namelist():
                                    # Get only the filename part for extension check and base name extraction
                                    member_filename_only = os.path.basename(member_name_in_zip)
                                    member_base_name, member_ext_zip_ext = os.path.splitext(member_filename_only)
                                    
                                    if member_name_in_zip.endswith('/') or member_name_in_zip.startswith('__MACOSX'): continue
                                    
                                    if member_ext_zip_ext.lower() in [".docx", ".pptx", ".pdf"]:
                                        members_found_for_processing_in_zip = True
                                        st.caption(f"  Processing member: {member_name_in_zip}")
                                        member_bytes = zip_ref.read(member_name_in_zip)
                                        member_stream = io.BytesIO(member_bytes)
                                        redacted_member_content = None
                                        
                                        if member_ext_zip_ext.lower() == ".docx": redacted_member_content = redact_docx(member_stream, names_variations_list, redaction_text_docx_pptx)
                                        elif member_ext_zip_ext.lower() == ".pptx": redacted_member_content = redact_pptx(member_stream, names_variations_list, redaction_text_docx_pptx)
                                        elif member_ext_zip_ext.lower() == ".pdf": redacted_member_content = redact_pdf(member_stream, names_variations_list, redaction_text_docx_pptx)
                                        
                                        if redacted_member_content:
                                            processed_zip_members_data.append((member_name_in_zip, redacted_member_content)) # Store full original path for structure
                                            st.success(f"    ✅ Redacted: {member_filename_only}")
                                            overall_docs_modified_count += 1
                                        else: 
                                            member_stream.seek(0)
                                            processed_zip_members_data.append((member_name_in_zip, member_stream))
                                            st.info(f"    ℹ️ No redactions in: {member_filename_only} (original included)")
                            
                            if members_found_for_processing_in_zip and processed_zip_members_data:
                                output_zip_stream = io.BytesIO()
                                file_counter_in_zip = 1
                                with zipfile.ZipFile(output_zip_stream, 'w', zipfile.ZIP_DEFLATED) as new_zip_archive:
                                    for m_full_path_in_zip, m_stream_content in processed_zip_members_data:
                                        m_stream_content.seek(0)
                                        
                                        original_member_filename_only = os.path.basename(m_full_path_in_zip)
                                        original_member_base_filename_only, original_member_ext_only = os.path.splitext(original_member_filename_only)
                                        
                                        extracted_initials = extract_initials_from_filename(original_member_base_filename_only)
                                        
                                        if extracted_initials: 
                                            new_filename_for_zip_entry_base = f"{actual_output_zip_base_name}_{extracted_initials}_{file_counter_in_zip:04d}"
                                        else: 
                                            new_filename_for_zip_entry_base = f"{actual_output_zip_base_name}_{file_counter_in_zip:04d}"
                                        
                                        new_filename_for_zip_entry = f"{new_filename_for_zip_entry_base}{original_member_ext_only}"
                                        
                                        # This ensures files are at the top level of the output ZIP
                                        new_zip_archive.writestr(new_filename_for_zip_entry, m_stream_content.read())
                                        file_counter_in_zip +=1
                                
                                output_zip_stream.seek(0)
                                display_zip_name = f"{actual_output_zip_base_name}.zip"
                                files_for_download.append((original_input_name, display_zip_name, output_zip_stream, ".zip"))
                                st.success(f"📦 Created new ZIP: **{display_zip_name}**.")
                                st.caption(f"   Files inside '{display_zip_name}' are at the top level, renamed using base '{actual_output_zip_base_name}', auto-extracted initials (if any), and a number.")

                            elif not members_found_for_processing_in_zip: 
                                st.info(f"ℹ️ No supported files (.docx, .pptx, .pdf) found in ZIP: **{original_input_name}** to process.")
                        except zipfile.BadZipFile: st.error(f"❌ Error: ZIP '{original_input_name}' appears to be corrupted.")
                        except Exception as e: st.error(f"❌ Error processing ZIP '{original_input_name}': {e}")
                    
                    else: # Process individual (non-ZIP) files
                        uploaded_file_obj.seek(0)
                        st.write(f"--- Processing: **{original_input_name}** ---")
                        redacted_content = None
                        try:
                            if file_extension == ".docx": redacted_content = redact_docx(uploaded_file_obj, names_variations_list, redaction_text_docx_pptx)
                            elif file_extension == ".pptx": redacted_content = redact_pptx(uploaded_file_obj, names_variations_list, redaction_text_docx_pptx)
                            elif file_extension == ".pdf": redacted_content = redact_pdf(uploaded_file_obj, names_variations_list, redaction_text_docx_pptx)
                            
                            if redacted_content:
                                display_name = f"redacted_{original_input_name}"
                                files_for_download.append((original_input_name, display_name, redacted_content, file_extension))
                                st.success(f"✅ Successfully redacted: **{original_input_name}**")
                                overall_docs_modified_count += 1
                            else:
                                st.info(f"ℹ️ No redactions made in **{original_input_name}** (file unchanged).")
                                uploaded_file_obj.seek(0) 
                                display_name = f"original_{original_input_name}"
                                files_for_download.append((original_input_name, display_name, uploaded_file_obj, file_extension))
                        except Exception as e: st.error(f"❌ Error processing **{original_input_name}**: {e}")

            if files_for_download:
                st.markdown("---"); st.subheader("⬇️ Download Files")
                num_files = len(files_for_download)
                max_cols = 3 
                num_cols_to_use = min(num_files, max_cols) if num_files > 0 else 1
                cols = st.columns(num_cols_to_use)
                for i, (orig_name, display_name, data_stream, ext) in enumerate(files_for_download):
                    mime_types = { ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                   ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                   ".pdf": "application/pdf", ".zip": "application/zip"}
                    mime_type = mime_types.get(ext, "application/octet-stream")
                    data_stream.seek(0) 
                    col_index = i % num_cols_to_use
                    with cols[col_index]:
                        st.download_button(
                            label=f"Download {display_name}", data=data_stream,
                            file_name=display_name, mime=mime_type,
                            key=f"dl_btn_{i}_{display_name.replace(' ','_').replace('.','_').replace('/','_')}"
                        )
            
            if overall_docs_modified_count == 0 and uploaded_files:
                st.info("ℹ️ No documents were modified based on the provided names, or no supported files were found in ZIPs.")
            elif overall_docs_modified_count > 0:
                st.balloons()

st.markdown("---")
with st.expander("📜 Instructions & Notes", expanded=False):
    st.markdown("""
        #### How to Use:
        1.  **Configure Settings (Sidebar):**
            *   Enter **Names to redact**: Provide a comma-separated list (e.g., `Professor Plum, John K. Doe, Ms. Scarlet, Young`). The app automatically generates variations (like `Plum` from `Professor Plum`, or `Mr. Young`, `Dr. Young` from `Young`) to improve detection.
            *   Set custom **Redaction text** for DOCX/PPTX files if desired.
            *   Optionally, for **ZIP Output Options** (these apply *only* when you upload a .zip file):
                *   **Custom output ZIP name base:** If you input `MyProject`, an uploaded `data.zip` will result in `MyProject.zip`. If left blank, it defaults to `redacted_data.zip`.
                *   Initials for renaming files *inside* the ZIP are automatically extracted from the original filenames of the member files (e.g., a file named `John-Doe-Report.docx` might contribute `JD` as initials). Max 2 initials.
        2.  **Upload Files (Main Area):** Select one or more .docx, .pptx, .pdf, or .zip files.
        3.  **Redact:** Click the "Redact Files" button.
        4.  **Download:**
            *   Processed individual files (or original if no changes) will be available for download.
            *   If you uploaded a ZIP file, a **new ZIP archive** will be created with the (potentially custom) name. All processed files from the original ZIP will be placed at the **top level** of this new ZIP, renamed as: `[OutputZipNameBase]_[ExtractedInitials]_[Number].ext`. If initials cannot be extracted, that part is omitted.

        #### Important Notes:
        *   **PDF Redaction:** Names in PDF files are "blacked out." The custom redaction text does not apply to PDF visual output.
        *   **ZIP File Processing:** Files within the original ZIP that are not .docx, .pptx, or .pdf are currently **not** included in the output ZIP.
        *   **Complex Documents:** For very complex layouts, embedded objects, or scanned (image-based) PDFs without OCR text, redaction might be incomplete. This tool works best with text-based documents. The redaction in DOCX/PPTX happens segment-by-segment (run-by-run), so names split across different formatting segments might not be fully caught if the individual parts don't match a variation.
        *   **Backup:** Always keep a backup of your original files before redacting!
        """)
